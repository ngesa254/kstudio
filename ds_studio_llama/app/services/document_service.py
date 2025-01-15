from typing import List, Dict, Any, Union
import logging
import os
import sys
import json
import shutil
import tempfile
import base64
import re
from pathlib import Path
from datetime import datetime
import io
import traceback

from llama_index.core import Document, VectorStoreIndex, Settings, SimpleDirectoryReader
from llama_index.core.node_parser import SentenceSplitter
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core.prompts import PromptTemplate
from pdf2image import convert_from_path
from PIL import Image
import pandas as pd
from pptx import Presentation

# Configure logging
log_file = "/Users/hmwangila/Documents/Documents/datascience_projects/2024/ds_studio_llama/ds_studio_llama/logs/agent_framework.log"
os.makedirs(os.path.dirname(log_file), exist_ok=True)

# Configure root logger
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)

# Get logger for this module
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

# Ensure llama_index logging is captured
llama_logger = logging.getLogger('llama_index')
llama_logger.setLevel(logging.INFO)

class DocumentService:
    def __init__(self, vector_store_path: str):
        logger.info(f"Initializing DocumentService with vector store path: {vector_store_path}")
        self.vector_store_path = vector_store_path
        
        # Disable LLM usage
        Settings.llm = None
        logger.info("LLM disabled for indexing")
        
        # Set up embedding model
        model_path = "/Users/hmwangila/Documents/Documents/datascience_projects/2024/ds_studio_llama/ds_studio_llama/all-MiniLM-L6-v2"
        self.embed_model = HuggingFaceEmbedding(
            model_name=model_path,
            trust_remote_code=True
        )
        Settings.embed_model = self.embed_model
        logger.info("Embedding model initialized")
        
        # Initial chunk sizes
        self.chunk_sizes = {
            "default": 2048,
            "spreadsheets": 2048,
            "presentations": 4096,
            "structured": 2048,
            "images": 2048
        }
        
        # Customize prompt template
        self.qa_prompt_tmpl = PromptTemplate(
            template=(
                "Here are relevant excerpts and/or images from the document:\n"
                "---------------------\n"
                "{context_str}\n"
                "---------------------\n"
                "Based on these excerpts and images, please provide a detailed answer to the following question. "
                "If the answer cannot be found in the provided content, say 'I don't have enough information to answer that question.'\n"
                "Question: {query_str}\n"
                "Answer: "
            )
        )
        logger.info("Prompt template configured")

    def _get_document_type(self, file_path: str) -> str:
        """Determine document type from file extension."""
        ext = os.path.splitext(file_path)[1].lower()
        for doc_type, extensions in self.get_supported_formats().items():
            if ext in extensions:
                logger.info(f"Document type determined: {doc_type} for file: {file_path}")
                return doc_type
        logger.error(f"Unsupported file format: {ext}")
        raise ValueError(f"Unsupported file format: {ext}")

    def _adjust_chunk_size(self, error_message: str, current_size: int) -> int:
        """Dynamically adjust chunk size based on error message."""
        try:
            # Extract metadata length from error message
            match = re.search(r"Metadata length \((\d+)\)", error_message)
            if match:
                metadata_length = int(match.group(1))
                # Add 20% buffer to the required size
                new_size = int(metadata_length * 1.2)
                logger.info(f"Adjusting chunk size from {current_size} to {new_size}")
                return new_size
        except Exception as e:
            logger.error(f"Error parsing metadata length: {str(e)}")
        
        # If parsing fails, double the current size
        new_size = current_size * 2
        logger.info(f"Doubling chunk size from {current_size} to {new_size}")
        return new_size

    def _process_spreadsheet(self, file_path: str) -> List[Document]:
        """Process spreadsheet files (Excel, CSV) using pandas."""
        logger.info(f"Processing spreadsheet: {file_path}")
        try:
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
                text = df.to_string(index=False)
                return [Document(text=text)]
            else:  # Excel files
                df = pd.read_excel(file_path, sheet_name=None)
                if isinstance(df, dict):  # Multiple sheets
                    documents = []
                    for sheet_name, sheet_df in df.items():
                        text = f"Sheet: {sheet_name}\n" + sheet_df.to_string(index=False)
                        documents.append(Document(text=text))
                    return documents
                else:
                    text = df.to_string(index=False)
                    return [Document(text=text)]
        except Exception as e:
            logger.error(f"Error processing spreadsheet: {str(e)}")
            raise

    def _process_presentation(self, file_path: str) -> List[Document]:
        """Process PowerPoint files using python-pptx."""
        logger.info(f"Processing presentation: {file_path}")
        try:
            # Extract text and notes using python-pptx
            prs = Presentation(file_path)
            text_content = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                # Get slide text
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                # Get slide notes
                notes = ""
                if slide.has_notes_slide and slide.notes_slide:
                    for shape in slide.notes_slide.shapes:
                        if hasattr(shape, "text"):
                            note_text = shape.text.strip()
                            if note_text:
                                notes += note_text + "\n"
                
                # Combine slide content
                slide_content = f"Slide {slide_number}:\n"
                if slide_text:
                    slide_content += "\nContent:\n" + "\n".join(slide_text)
                if notes:
                    slide_content += "\nNotes:\n" + notes
                
                text_content.append(slide_content)
            
            # Create document with text
            full_text = "\n\n".join(text_content)
            return [Document(text=full_text)]
            
        except Exception as e:
            logger.error(f"Error processing presentation: {str(e)}\n{traceback.format_exc()}")
            raise

    def _process_image(self, file_path: str) -> List[Document]:
        """Process image files for Claude."""
        logger.info(f"Processing image file: {file_path}")
        try:
            # Open and resize image
            with Image.open(file_path) as img:
                # Calculate new dimensions while maintaining aspect ratio
                max_size = (300, 300)  # Maximum width/height
                ratio = min(max_size[0] / img.width, max_size[1] / img.height)
                new_size = (int(img.width * ratio), int(img.height * ratio))
                
                # Resize image
                img = img.resize(new_size, Image.Resampling.LANCZOS)
                
                # Convert to RGB if necessary
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                
                # Convert to base64
                buffered = io.BytesIO()
                img.save(buffered, format="JPEG", quality=85)
                base64_img = base64.b64encode(buffered.getvalue()).decode()
            
            metadata = {
                "images": [{
                    "type": "image",
                    "data": base64_img
                }]
            }
            
            # Let Claude interpret the image
            text = "Please analyze this image and answer any questions about its content."
            
            return [Document(text=text, metadata=metadata)]
        except Exception as e:
            logger.error(f"Error processing image: {str(e)}")
            raise

    def _get_reader_kwargs(self, doc_type: str, file_path: str = None) -> Dict[str, Any]:
        """Get specialized reader arguments based on document type."""
        logger.info(f"Getting reader configuration for document type: {doc_type}")
        base_kwargs = {
            "recursive": False,
            "filename_as_id": True,
            "num_files_limit": 1
        }
        
        if doc_type == "spreadsheets":
            # Handle spreadsheets directly through pandas
            if file_path:
                documents = self._process_spreadsheet(file_path)
                base_kwargs["documents"] = documents
            logger.info("Using spreadsheet configuration with pandas")
        elif doc_type == "presentations":
            # Handle presentations directly
            if file_path:
                documents = self._process_presentation(file_path)
                base_kwargs["documents"] = documents
            logger.info("Using presentation configuration")
        elif doc_type == "images":
            # Handle images directly
            if file_path:
                documents = self._process_image(file_path)
                base_kwargs["documents"] = documents
            logger.info("Using image configuration")
        elif doc_type == "structured":
            base_kwargs.update({
                "encoding": "utf-8",
                "include_metadata": True
            })
            logger.info("Using structured document configuration")
        
        return base_kwargs

    def _get_node_parser(self, doc_type: str) -> SentenceSplitter:
        """Get specialized node parser based on document type."""
        logger.info(f"Getting node parser for document type: {doc_type}")
        chunk_size = self.chunk_sizes.get(doc_type, self.chunk_sizes["default"])
        overlap = int(chunk_size * 0.1)  # 10% overlap
        
        parser = SentenceSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            paragraph_separator="\n\n"
        )
        logger.info(f"Using node parser with chunk size {chunk_size} and overlap {overlap}")
        return parser

    async def process_document(self, file_path: str, collection_name: str) -> bool:
        """Process a document and store it in the vector store."""
        logger.info(f"Processing document: {file_path} for collection: {collection_name}")
        temp_dir = None
        try:
            if not os.path.exists(file_path):
                logger.error(f"Document not found: {file_path}")
                raise FileNotFoundError(f"Document not found at path: {file_path}")

            # Determine document type
            doc_type = self._get_document_type(file_path)
            logger.info(f"Processing {doc_type} document: {file_path}")
            
            # Create temporary directory for processing
            temp_dir = tempfile.mkdtemp()
            temp_file = os.path.join(temp_dir, os.path.basename(file_path))
            shutil.copy2(file_path, temp_file)
            logger.info(f"Copied file to temp location: {temp_file}")

            # Get reader configuration
            reader_kwargs = self._get_reader_kwargs(doc_type, file_path)  # Use original file path
            
            # Get documents either from direct processing or SimpleDirectoryReader
            if "documents" in reader_kwargs:
                documents = reader_kwargs.pop("documents")
                logger.info("Using pre-processed documents")
            else:
                # Use SimpleDirectoryReader for other document types
                logger.info(f"Using SimpleDirectoryReader with {doc_type} configuration")
                reader = SimpleDirectoryReader(
                    input_dir=temp_dir,
                    **reader_kwargs
                )
                documents = reader.load_data()
            
            if not documents:
                logger.warning("No content found in document")
                raise ValueError("No content found in document")
            
            logger.info(f"Loaded document with {len(documents)} sections")
            
            # Create vector store directory if it doesn't exist
            collection_path = os.path.join(self.vector_store_path, collection_name)
            os.makedirs(collection_path, exist_ok=True)
            
            while True:
                try:
                    node_parser = self._get_node_parser(doc_type)
                    logger.info(f"Using specialized node parser for {doc_type}")
                    
                    logger.info("Creating vector store index")
                    index = VectorStoreIndex.from_documents(
                        documents=documents,
                        transformations=[node_parser],
                        show_progress=True
                    )
                    
                    logger.info(f"Saving index to: {collection_path}")
                    index.storage_context.persist(persist_dir=collection_path)
                    logger.info("Successfully saved index")
                    
                    return True
                    
                except Exception as e:
                    error_msg = str(e)
                    if "Metadata length" in error_msg and "longer than chunk size" in error_msg:
                        # Get current chunk size and adjust it
                        current_size = self.chunk_sizes.get(doc_type, self.chunk_sizes["default"])
                        new_size = self._adjust_chunk_size(error_msg, current_size)
                        
                        # Update chunk size for this document type
                        self.chunk_sizes[doc_type] = new_size
                        logger.info(f"Retrying with new chunk size: {new_size}")
                        continue
                    
                    logger.error(f"Error creating index: {error_msg}\n{traceback.format_exc()}")
                    if os.path.exists(collection_path):
                        shutil.rmtree(collection_path)
                    raise ValueError(f"Failed to create index: {error_msg}")
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}\n{traceback.format_exc()}")
            raise
        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
                logger.info("Cleaned up temporary directory")

    async def delete_collection(self, collection_name: str) -> bool:
        """Delete a collection."""
        logger.info(f"Deleting collection: {collection_name}")
        try:
            collection_path = os.path.join(self.vector_store_path, collection_name)
            if os.path.exists(collection_path):
                shutil.rmtree(collection_path)
                logger.info(f"Deleted collection: {collection_name}")
            return True
        except Exception as e:
            logger.error(f"Error deleting collection: {str(e)}")
            raise

    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get list of supported formats with their categories."""
        formats = {
            "documents": [".pdf", ".txt", ".doc", ".docx"],
            "spreadsheets": [".xlsx", ".xls", ".csv"],
            "presentations": [".pptx", ".ppt"],
            "structured": [".json", ".md", ".html"],
            "images": [".png", ".jpg", ".jpeg", ".gif", ".bmp"]
        }
        logger.info(f"Supported formats: {formats}")
        return formats

    async def get_document_context(self, collection_name: str, query: str = None) -> List[Dict[str, Any]]:
        """Retrieve context from a document collection.
        If query is provided, returns relevant context for the query.
        Otherwise returns full document context."""
        logger.info(f"Getting document context for collection: {collection_name}, query: {query}")
        try:
            collection_path = os.path.join(self.vector_store_path, collection_name)
            if not os.path.exists(collection_path):
                logger.error(f"Collection not found: {collection_name}")
                raise ValueError(f"Collection not found: {collection_name}")

            logger.info("Loading index")
            index = VectorStoreIndex.load_from_disk(
                persist_dir=collection_path,
                embed_model=self.embed_model
            )

            if query:
                logger.info(f"Executing query: {query}")
                query_engine = index.as_query_engine(
                    similarity_top_k=5,
                    response_mode="no_text"
                )
                response = query_engine.query(query)
                context = [{
                    "text": node.text,
                    "score": node.score,
                    "images": node.metadata.get("images", []) if hasattr(node, "metadata") else []
                } for node in response.source_nodes]
                logger.info(f"Found {len(context)} relevant context items")
                return context
            else:
                logger.info("Getting full document context")
                all_nodes = index.docstore.docs.values()
                context = [{
                    "text": node.text,
                    "id": node.doc_id,
                    "images": node.metadata.get("images", []) if hasattr(node, "metadata") else []
                } for node in all_nodes]
                logger.info(f"Retrieved {len(context)} context items")
                return context

        except Exception as e:
            logger.error(f"Error retrieving context: {str(e)}")
            raise
